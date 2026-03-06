#!/usr/bin/env python3
"""
Exporta os dados do dashboard para uma apresentação (PPTX) por mês.

Uso:
  python export_pptx.py               # usa mês atual, tenta data/meses/YYYY-MM.json
  python export_pptx.py 2026-03       # exporta mês específico
"""

from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DATA_DIR = ROOT / "data"
DOCS_DIR = ROOT / "docs"


def _mes_key(dt: datetime) -> str:
    return dt.strftime("%Y-%m")


def _load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def _safe_get(info: dict, *path: str, default: str = "—") -> str:
    cur: Any = info
    for p in path:
        if not isinstance(cur, dict) or p not in cur:
            return default
        cur = cur[p]
    return str(cur) if cur is not None else default


def export_mes(mes: str) -> Path:
    month_path = DATA_DIR / "meses" / f"{mes}.json"
    if month_path.exists():
        payload = _load_json(month_path)
        dados = payload.get("dados", payload)
        atualizado_em = payload.get("atualizado_em") or ""
    else:
        dados = _load_json(DATA_DIR / "metas_atual.json")
        atualizado_em = ""

    prs = Presentation()

    # Slide 1: capa
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(12.0), Inches(1.0))
    tf = title_box.text_frame
    tf.text = "Dashboard de Metas — Top Estética Bucal"
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(12.0), Inches(0.8))
    stf = sub_box.text_frame
    stf.text = f"Mês: {mes}"
    stf.paragraphs[0].font.size = Pt(22)

    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(12.0), Inches(0.8))
    itf = info_box.text_frame
    itf.text = f"Gerado em: {datetime.now().strftime('%d/%m/%Y %H:%M')}" + (f" | Atualizado em: {atualizado_em}" if atualizado_em else "")
    itf.paragraphs[0].font.size = Pt(14)

    # Slides por cidade
    for cidade in sorted(dados.keys(), key=lambda s: s.lower()):
        info = dados[cidade] or {}
        s = prs.slides.add_slide(prs.slide_layouts[5])

        h = s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.7))
        ht = h.text_frame
        ht.text = cidade
        ht.paragraphs[0].font.size = Pt(30)
        ht.paragraphs[0].font.bold = True

        box = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12.0), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True

        lines = [
            f"Ortodontia: {_safe_get(info, 'metas_financeiras', 'ortodontia')}",
            f"Clínico Geral: {_safe_get(info, 'metas_financeiras', 'clinico_geral')}",
            f"Profilaxia: {_safe_get(info, 'metas_servicos', 'profilaxia')}",
            f"Restauração: {_safe_get(info, 'metas_servicos', 'restauracao')}",
            f"Novas avaliações Google: {_safe_get(info, 'avaliacoes', 'novas', default='0')}",
            f"Total avaliações (atual): {_safe_get(info, 'avaliacoes', 'atual', default='0')}",
        ]

        tf.text = lines[0]
        tf.paragraphs[0].font.size = Pt(18)
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)

        footer = s.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(12.0), Inches(0.4))
        ft = footer.text_frame
        ft.text = "Top Estética Bucal — Export mensal"
        ft.paragraphs[0].font.size = Pt(12)
        ft.paragraphs[0].alignment = PP_ALIGN.RIGHT

    out_dir = DATA_DIR / "exports" / mes
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y-%m-%d_%H%M")
    out_path = out_dir / f"dashboard_metas_{mes}_{stamp}.pptx"
    prs.save(out_path)

    # Copia para docs/exports para download no dashboard (se estiver hospedado)
    docs_out_dir = DOCS_DIR / "exports" / mes
    docs_out_dir.mkdir(parents=True, exist_ok=True)
    docs_copy = docs_out_dir / out_path.name
    docs_copy.write_bytes(out_path.read_bytes())

    # Nome estável "latest" para facilitar link no dashboard
    (out_dir / f"dashboard_metas_{mes}_latest.pptx").write_bytes(out_path.read_bytes())
    (docs_out_dir / f"dashboard_metas_{mes}_latest.pptx").write_bytes(out_path.read_bytes())

    return out_path


def main() -> None:
    mes = sys.argv[1].strip() if len(sys.argv) > 1 else _mes_key(datetime.now())
    out = export_mes(mes)
    print(f"PPTX gerado: {out}")


if __name__ == "__main__":
    main()

